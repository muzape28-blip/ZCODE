# python-pptx — membuat presentasi PowerPoint
# Butuh: install python-pptx
# Hasil: presentasi_zcode.pptx di workspace

from pptx import Presentation

presentasi = Presentation()

judul = presentasi.slides.add_slide(presentasi.slide_layouts[0])
judul.shapes.title.text = "Belajar Python dari HP"
judul.placeholders[1].text = "Presentasi dibuat otomatis dengan ZCODE"

isi = presentasi.slides.add_slide(presentasi.slide_layouts[1])
isi.shapes.title.text = "Yang sudah dipelajari"
frame = isi.placeholders[1].text_frame
frame.text = "Dasar Python"
for teks in ["Mengolah data", "Membuat file Office", "Mengakses API"]:
    frame.add_paragraph().text = teks

presentasi.save("presentasi_zcode.pptx")
print("Tersimpan: presentasi_zcode.pptx")
